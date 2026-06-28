import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Use blank slide layout
    blank_layout = prs.slide_layouts[6]
    
    brand_bg = RGBColor(11, 15, 25)
    brand_text = RGBColor(255, 255, 255)
    brand_accent = RGBColor(0, 230, 118) # Neon Green
    brand_danger = RGBColor(255, 71, 87) # Coral Red
    
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = brand_bg
        
    def add_title(slide, text):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Century Gothic'
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = brand_accent
        
    def add_bullet(tf, text, level=0, bold=False, color=brand_text):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = 'Segoe UI'
        p.font.size = Pt(22)
        p.font.color.rgb = color
        p.font.bold = bold

    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    
    poster_path = r"C:\Users\Admin\.gemini\antigravity\brain\d7d9d61b-485b-428c-ab2a-bb6514e780b2\saferoute_twilight_poster_1781724977743.png"
    if os.path.exists(poster_path):
        slide1.shapes.add_picture(poster_path, Inches(0), Inches(0), Inches(10), Inches(5.625))
    
    txBox1 = slide1.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.5))
    tf1 = txBox1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "SafeRoute"
    p1.font.name = 'Century Gothic'
    p1.font.size = Pt(60)
    p1.font.bold = True
    p1.font.color.rgb = brand_accent
    
    p2 = tf1.add_paragraph()
    p2.text = "The AI-Powered Sentinel Network"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(24)
    p2.font.color.rgb = brand_text

    # Slide 2: The Problem
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_title(slide2, "The Problem")
    
    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
    tf2 = txBox2.text_frame
    add_bullet(tf2, "Current navigation apps optimize for speed, completely ignoring safety.")
    add_bullet(tf2, "In emergencies, unlocking a phone and opening an app takes too long.")
    add_bullet(tf2, "Users have no real-time visibility into crowd behaviors or unlit streets.")

    # Slide 3: The Solution
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_title(slide3, "The Solution")
    
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
    tf3 = txBox3.text_frame
    add_bullet(tf3, "AI Threat Routing:", bold=True, color=brand_danger)
    add_bullet(tf3, "Re-calculates routes dynamically to avoid known 'Threat Zones'.", level=1)
    add_bullet(tf3, "The Sentinel Network:", bold=True, color=brand_danger)
    add_bullet(tf3, "A decentralized mesh of trusted contacts and nearby users.", level=1)
    add_bullet(tf3, "ARIA (AI Guardian):", bold=True, color=brand_danger)
    add_bullet(tf3, "A built-in AI assistant that pre-scans zones for danger.", level=1)

    # Slide 4: The Tech Stack
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_title(slide4, "The Tech Stack")
    
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    tf4 = txBox4.text_frame
    add_bullet(tf4, "Frontend")
    add_bullet(tf4, "React + Vite", level=1)
    add_bullet(tf4, "Tailwind CSS (Glassmorphism)", level=1)
    
    txBox4_b = slide4.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
    tf4_b = txBox4_b.text_frame
    add_bullet(tf4_b, "Backend & Integrations")
    add_bullet(tf4_b, "Python FastAPI + SQLite", level=1)
    add_bullet(tf4_b, "Twilio REST API (Instant SOS)", level=1)
    add_bullet(tf4_b, "Leaflet & OSRM (Dynamic Maps)", level=1)

    # Slide 5: The Live Demo
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_title(slide5, "Live Demo")
    
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
    tf5 = txBox5.text_frame
    add_bullet(tf5, "1. Plot a Route", bold=True)
    add_bullet(tf5, "Showcase dynamic navigation and threat avoidance.", level=1)
    add_bullet(tf5, "2. Trigger Guardian Pulse", bold=True)
    add_bullet(tf5, "Use the instant double-tap feature.", level=1)
    add_bullet(tf5, "3. Real-World Alert", bold=True)
    add_bullet(tf5, "Receive physical Twilio SMS with live coordinates instantly.", level=1)

    prs.save('SafeRoute_Presentation.pptx')
    print("Presentation created successfully.")

if __name__ == '__main__':
    create_presentation()
