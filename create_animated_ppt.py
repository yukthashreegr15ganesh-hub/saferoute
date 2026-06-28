import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml

def add_transition(slide):
    # Adds a smooth fade transition to the slide
    transition_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow"><p:fade/></p:transition>'
    try:
        slide.element.append(parse_xml(transition_xml))
    except:
        pass

def create_presentation():
    prs = Presentation()
    # Widescreen 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Theme colors based on Sapthagiri NPS University Logo
    navy_blue = RGBColor(0, 32, 96)
    gold = RGBColor(218, 165, 32)
    white = RGBColor(255, 255, 255)
    light_gray = RGBColor(240, 240, 240)
    
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    def add_title(slide, text):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Impact' # Very attractive bold font
        p.font.size = Pt(54)
        p.font.color.rgb = gold
        p.alignment = PP_ALIGN.LEFT

    def add_subtitle(slide, text, top):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(32)
        p.font.color.rgb = white
        p.font.italic = True
        p.alignment = PP_ALIGN.LEFT

    def add_content(slide, title, bullets, image_path=None):
        set_slide_bg(slide, navy_blue)
        add_transition(slide)
        add_title(slide, title)
        
        # Add decorative line
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.6), Inches(12.333), Inches(0.05) # 1 is msoShapeRectangle
        )
        line.fill.solid()
        line.fill.fore_color.rgb = gold
        line.line.color.rgb = gold

        width = Inches(12.333) if not image_path else Inches(7)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), width, Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, (text, level) in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = text
            p.level = level
            p.font.name = 'Arial' if level == 0 else 'Calibri'
            p.font.size = Pt(32 if level == 0 else 24)
            p.font.color.rgb = white if level == 0 else light_gray
            p.font.bold = (level == 0)

        if image_path and os.path.exists(image_path):
            try:
                slide.shapes.add_picture(image_path, Inches(7.5), Inches(2.2), width=Inches(5.333))
            except:
                pass

    # Media paths
    media1 = r"C:\Users\Admin\.gemini\antigravity\brain\d7d9d61b-485b-428c-ab2a-bb6514e780b2\media__1781725715177.png"
    media2 = r"C:\Users\Admin\.gemini\antigravity\brain\d7d9d61b-485b-428c-ab2a-bb6514e780b2\media__1781725750203.png"
    poster = r"C:\Users\Admin\.gemini\antigravity\brain\d7d9d61b-485b-428c-ab2a-bb6514e780b2\saferoute_isometric_poster_1781725101832.png"

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, navy_blue)
    add_transition(slide1)
    
    # Building as background or large picture
    if os.path.exists(media2):
        try:
            slide1.shapes.add_picture(media2, Inches(0), Inches(0), width=Inches(13.333))
            # Add a semi-transparent overlay to make text readable (using a dark rectangle)
            overlay = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = navy_blue
            # python-pptx doesn't support setting transparency easily, so we just use the shape if needed
            # Actually we'll just put the building on the right side if we can't do transparency easily
        except:
            pass
            
    # Redo Slide 1 to make it perfect
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, navy_blue)
    add_transition(slide1)

    # Logo
    if os.path.exists(media1):
        try:
            slide1.shapes.add_picture(media1, Inches(0.5), Inches(0.5), height=Inches(1.5))
        except:
            pass

    # Building Picture (Attractive placement)
    if os.path.exists(media2):
        try:
            slide1.shapes.add_picture(media2, Inches(6.5), Inches(0), height=Inches(7.5))
        except:
            pass

    # Title Text
    txBox1 = slide1.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(6), Inches(2))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "SafeRoute"
    p1.font.name = 'Impact'
    p1.font.size = Pt(80)
    p1.font.color.rgb = gold
    
    p2 = tf1.add_paragraph()
    p2.text = "The AI-Powered Sentinel Network"
    p2.font.name = 'Trebuchet MS'
    p2.font.size = Pt(28)
    p2.font.color.rgb = white
    p2.font.italic = True

    # --- Slide 2: Problem ---
    add_content(prs.slides.add_slide(blank_layout), "THE CHALLENGE", [
        ("Current navigation apps lack safety focus", 0),
        ("They optimize for speed, completely ignoring unlit streets or danger zones.", 1),
        ("Emergency response is too slow", 0),
        ("Unlocking a phone, opening an app, and dialing takes precious seconds.", 1),
        ("No real-time crowd intelligence", 0),
        ("Users cannot see live threat zones or dispersing crowds ahead of time.", 1)
    ])

    # --- Slide 3: The Solution ---
    add_content(prs.slides.add_slide(blank_layout), "OUR SOLUTION", [
        ("AI Threat Routing", 0),
        ("Dynamically re-calculates paths to avoid identified 'Threat Zones'.", 1),
        ("The Sentinel Network", 0),
        ("A decentralized mesh of trusted contacts providing live updates.", 1),
        ("ARIA (AI Guardian)", 0),
        ("A built-in AI assistant that pre-scans zones for proactive defense.", 1)
    ], poster)

    # --- Slide 4: Technology Stack ---
    add_content(prs.slides.add_slide(blank_layout), "TECHNOLOGY STACK", [
        ("Frontend Application", 0),
        ("React + Vite with advanced Glassmorphism UI.", 1),
        ("Backend Systems", 0),
        ("Python FastAPI + SQLite for high-speed data processing.", 1),
        ("Critical Integrations", 0),
        ("Twilio REST API for sub-second SOS SMS broadcasting.", 1),
        ("Leaflet & OSRM for dynamic mapping.", 1)
    ])

    # --- Slide 5: Live Demonstration ---
    add_content(prs.slides.add_slide(blank_layout), "LIVE DEMONSTRATION", [
        ("1. Plot a Route", 0),
        ("Showcase dynamic navigation and threat avoidance algorithms.", 1),
        ("2. Trigger Guardian Pulse", 0),
        ("Use the secret instant double-tap feature.", 1),
        ("3. Real-World Alert", 0),
        ("Receive physical Twilio SMS with live coordinates instantly.", 1)
    ])

    # Remove the blank first slide we abandoned
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    prs.save('SafeRoute_Final_Presentation.pptx')
    print("Beautiful presentation created successfully.")

if __name__ == '__main__':
    create_presentation()
